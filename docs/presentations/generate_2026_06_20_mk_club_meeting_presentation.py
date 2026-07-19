"""Build the 2026-06-20 Mike & Key club presentation slide deck."""

from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
import re
from typing import Iterable

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt


BASE_DIR = Path(__file__).resolve().parent
IMAGE_DIR = BASE_DIR / 'images'
OUTPUT_PATH = (
    BASE_DIR / '2026-06-20-mk-club-meeting-presentation.pptx'
)

SLIDE_WIDTH = Cm(33.87)
SLIDE_HEIGHT = Cm(19.05)
TITLE_BAR_HEIGHT = Cm(2.5)
FULL_BLEED_BAR_HEIGHT = Cm(2.6)
SECTION_BAR_HEIGHT = Cm(4.0)
MARGIN = Cm(0.9)
SMALL_LOGO = Cm(1.8)

FONT_NAME = 'Liberation Sans'
FALLBACK_FONT = 'Arial'
TITLE_SIZE = Pt(36)
SECTION_NUMBER_SIZE = Pt(28)
BODY_SIZE = Pt(22)
CAPTION_SIZE = Pt(14)
CTA_SIZE = Pt(24)

BG = '1A1A1A'
TITLE_BG = '1B4D18'
ACCENT = '2D7A27'
LIGHT_ACCENT = '4CAF50'
BODY = 'F0F4F8'
SECONDARY = 'A8A8A8'
AMBER = 'F0A500'

FULL_SLIDE_PX = (1920, 1080)


@dataclass(frozen=True)
class ImagePlacement:
    """Placement specification for an individual slide image."""

    image: str
    x: int
    y: int
    w: int
    h: int
    mode: str = 'cover'
    placeholder_text: str = ''


@dataclass(frozen=True)
class SlideSpec:
    """Content specification for a single slide."""

    layout: str
    title: str = ''
    subtitle: str = ''
    presenter: str = ''
    section_number: str = ''
    image: str = ''
    image_caption: str = ''
    bullets: tuple[str, ...] = ()
    cta: str = ''
    notes: tuple[str, ...] = ()
    body_image_mode: str = 'cover'
    placeholder_text: str = ''
    body_font_size: int | None = None
    bullet_rect: tuple[int, int, int, int] | None = None
    image_rect: tuple[int, int, int, int] | None = None
    extra_images: tuple[ImagePlacement, ...] = ()
    mosaic_images: tuple[str, ...] = ()
    qr_image: str = ''
    qr_rect: tuple[int, int, int, int] | None = None
    show_logo: bool = True


def rgb(value: str) -> RGBColor:
    """Convert a hex color string into a PowerPoint RGB color."""

    return RGBColor.from_string(value)


def px_size(width: int, height: int) -> tuple[int, int]:
    """Scale a slide region into a reasonable raster target size."""

    slide_w = int(SLIDE_WIDTH)
    slide_h = int(SLIDE_HEIGHT)
    px_w = max(100, round(width / slide_w * FULL_SLIDE_PX[0]))
    px_h = max(100, round(height / slide_h * FULL_SLIDE_PX[1]))
    return px_w, px_h


def parse_markup(text: str) -> list[tuple[str, dict[str, bool]]]:
    """Split simple markdown-style emphasis into text fragments."""

    token_re = re.compile(r'(\*\*.*?\*\*|\*.*?\*)')
    parts: list[tuple[str, dict[str, bool]]] = []
    for chunk in token_re.split(text):
        if not chunk:
            continue
        if chunk.startswith('**') and chunk.endswith('**'):
            parts.append((chunk[2:-2], {'bold': True, 'italic': False}))
        elif chunk.startswith('*') and chunk.endswith('*'):
            parts.append((chunk[1:-1], {'bold': False, 'italic': True}))
        else:
            parts.append((chunk, {'bold': False, 'italic': False}))
    return parts


def add_background(slide) -> None:
    """Apply the standard dark presentation background."""

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(BG)


def add_rect(slide, x: int, y: int, w: int, h: int, color: str,
             transparency: float = 0.0):
    """Add a filled rectangle with no visible outline."""

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        x,
        y,
        w,
        h,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency
    shape.line.fill.background()
    return shape


def add_textbox(
    slide,
    x: int,
    y: int,
    w: int,
    h: int,
    text: str = '',
    *,
    font_size=BODY_SIZE,
    color: str = BODY,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_VERTICAL_ANCHOR.TOP,
):
    """Create a textbox with a single paragraph of styled text."""

    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = valign
    para = frame.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = font_size
    run.font.color.rgb = rgb(color)
    run.font.bold = bold
    run.font.italic = italic
    if text:
        run.font.name = FONT_NAME
    return box


def add_markup_paragraph(
    frame,
    text: str,
    *,
    color: str = BODY,
    font_size=BODY_SIZE,
    align=PP_ALIGN.LEFT,
    bullet: bool = False,
) -> None:
    """Append a paragraph with optional green manual bullet styling."""

    para = frame.paragraphs[0] if not frame.text else frame.add_paragraph()
    para.alignment = align
    para.space_after = Pt(10)
    para.line_spacing = 1.15
    if bullet:
        bullet_run = para.add_run()
        bullet_run.text = '• '
        bullet_run.font.name = FONT_NAME
        bullet_run.font.size = font_size
        bullet_run.font.bold = True
        bullet_run.font.color.rgb = rgb(ACCENT)
    for segment, style in parse_markup(text):
        run = para.add_run()
        run.text = segment
        run.font.name = FONT_NAME
        run.font.size = font_size
        run.font.color.rgb = rgb(color)
        run.font.bold = style['bold']
        run.font.italic = style['italic']


def add_bullets(
    slide,
    bullets: Iterable[str],
    x: int,
    y: int,
    w: int,
    h: int,
    *,
    font_size=BODY_SIZE,
) -> None:
    """Create a standard bullet list textbox."""

    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    for idx, bullet in enumerate(bullets):
        if idx == 0:
            frame.clear()
        text_color = AMBER if bullet.startswith('⚠️') else BODY
        add_markup_paragraph(
            frame,
            bullet,
            color=text_color,
            font_size=font_size,
            bullet=True,
        )


def safe_image_bytes(
    image_name: str,
    target_w: int,
    target_h: int,
    *,
    mode: str = 'cover',
) -> BytesIO | None:
    """Load, crop, and encode an image into a PNG byte stream."""

    path = IMAGE_DIR / image_name
    if not path.exists():
        return None
    try:
        with Image.open(path) as source:
            image = ImageOps.exif_transpose(source).convert('RGBA')
            target_size = px_size(target_w, target_h)
            if mode == 'cover':
                rendered = ImageOps.fit(
                    image,
                    target_size,
                    method=Image.Resampling.LANCZOS,
                    centering=(0.5, 0.5),
                )
            else:
                rendered = Image.new('RGBA', target_size, f'#{BG}')
                contain = ImageOps.contain(
                    image,
                    target_size,
                    method=Image.Resampling.LANCZOS,
                )
                left = (target_size[0] - contain.width) // 2
                top = (target_size[1] - contain.height) // 2
                rendered.paste(contain, (left, top), contain)
            stream = BytesIO()
            rendered.save(stream, format='PNG')
            stream.seek(0)
            return stream
    except Exception:
        return None


def add_placeholder(
    slide,
    x: int,
    y: int,
    w: int,
    h: int,
    text: str,
) -> None:
    """Add a dark placeholder when an image is missing or unreadable."""

    shape = add_rect(slide, x, y, w, h, BG)
    shape.line.fill.solid()
    shape.line.fill.fore_color.rgb = rgb(SECONDARY)
    shape.line.width = Pt(1.25)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    para = frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = rgb(SECONDARY)


def add_image(
    slide,
    image_name: str,
    x: int,
    y: int,
    w: int,
    h: int,
    *,
    mode: str = 'cover',
    placeholder_text: str = '',
) -> None:
    """Place an image or a readable placeholder in the given region."""

    stream = safe_image_bytes(image_name, w, h, mode=mode)
    if stream is None:
        label = placeholder_text or image_name
        add_placeholder(slide, x, y, w, h, label)
        return
    slide.shapes.add_picture(stream, x, y, width=w, height=h)


def add_image_placement(slide, placement: ImagePlacement) -> None:
    """Add a positioned image defined by an ImagePlacement spec."""

    add_image(
        slide,
        placement.image,
        placement.x,
        placement.y,
        placement.w,
        placement.h,
        mode=placement.mode,
        placeholder_text=placement.placeholder_text,
    )


def add_logo(
    slide,
    *,
    centered: bool = False,
    size=None,
    y=None,
) -> None:
    """Add the Mike & Key logo in the standard location."""

    logo_path = IMAGE_DIR / 'MnK_Logo.png'
    if not logo_path.exists():
        return
    if centered:
        size = size or Cm(6.0)
        x = (SLIDE_WIDTH - size) // 2
        y = Cm(1.8) if y is None else y
        slide.shapes.add_picture(str(logo_path), x, y, height=size)
        return
    x = int(SLIDE_WIDTH) - int(MARGIN) - int(SMALL_LOGO)
    y = int(SLIDE_HEIGHT) - int(MARGIN) - int(SMALL_LOGO)
    slide.shapes.add_picture(str(logo_path), x, y, height=SMALL_LOGO)


def add_title_bar(slide, title: str) -> None:
    """Add the branded title bar used on standard content slides."""

    add_rect(slide, 0, 0, int(SLIDE_WIDTH), int(TITLE_BAR_HEIGHT), TITLE_BG)
    box = slide.shapes.add_textbox(
        int(MARGIN),
        Cm(0.35),
        int(SLIDE_WIDTH) - int(Cm(2.8)),
        int(Cm(1.8)),
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    para = frame.paragraphs[0]
    for segment, style in parse_markup(title):
        run = para.add_run()
        run.text = segment
        run.font.name = FONT_NAME
        run.font.size = TITLE_SIZE
        run.font.bold = True or style['bold']
        run.font.italic = style['italic']
        run.font.color.rgb = rgb(BODY)


def add_notes(slide, lines: Iterable[str]) -> None:
    """Populate the presenter notes text frame for a slide."""

    frame = slide.notes_slide.notes_text_frame
    frame.clear()
    for idx, line in enumerate(lines):
        para = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        para.text = line


def add_optional_qr(slide, spec: SlideSpec) -> None:
    """Add an optional QR code image for slides that request one."""

    if not spec.qr_image or spec.qr_rect is None:
        return
    qr_x, qr_y, qr_w, qr_h = spec.qr_rect
    add_image(
        slide,
        spec.qr_image,
        qr_x,
        qr_y,
        qr_w,
        qr_h,
        mode='contain',
        placeholder_text='QR code',
    )


def add_mosaic_panel(slide, images: tuple[str, ...]) -> None:
    """Render a simple three-image mosaic on the left side of a slide."""

    panel_width = int(int(SLIDE_WIDTH) * 0.56)
    gutter = int(Cm(0.12))
    panel_height = int(SLIDE_HEIGHT)
    tile_height = (panel_height - (gutter * (len(images) - 1))) // len(images)

    y = 0
    for index, image_name in enumerate(images):
        height = (
            tile_height
            if index < len(images) - 1
            else panel_height - y
        )
        add_image(
            slide,
            image_name,
            0,
            y,
            panel_width,
            height,
            mode='cover',
            placeholder_text=image_name,
        )
        y += height + gutter


def finalize_slide(slide, spec: SlideSpec):
    """Add common per-slide finishing elements."""

    if spec.show_logo:
        add_logo(slide)
    add_optional_qr(slide, spec)
    add_notes(slide, spec.notes)
    return slide


def build_layout_a(prs: Presentation, spec: SlideSpec, closing: bool = False):
    """Build a title or closing slide."""

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    if closing:
        add_logo(slide, centered=True)
    else:
        add_logo(slide, centered=True, size=Cm(7.0), y=Cm(1.1))
    if closing:
        add_textbox(
            slide,
            Cm(6.0),
            Cm(8.5),
            Cm(21.8),
            Cm(2.4),
            'Questions?',
            font_size=Pt(34),
            color=BODY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            Cm(5.2),
            Cm(11.0),
            Cm(23.5),
            Cm(1.2),
            spec.subtitle,
            font_size=Pt(20),
            color=SECONDARY,
            align=PP_ALIGN.CENTER,
        )
        add_optional_qr(slide, spec)
        add_notes(slide, spec.notes)
        return slide

    add_textbox(
        slide,
        Cm(3.2),
        Cm(10.0),
        Cm(27.5),
        Cm(2.4),
        spec.title,
        font_size=TITLE_SIZE,
        color=BODY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Cm(5.4),
        Cm(12.7),
        Cm(23.0),
        Cm(1.3),
        spec.subtitle,
        font_size=Pt(24),
        color=LIGHT_ACCENT,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Cm(7.3),
        Cm(15.0),
        Cm(19.5),
        Cm(1.0),
        spec.presenter,
        font_size=Pt(20),
        color=SECONDARY,
        align=PP_ALIGN.CENTER,
    )
    add_notes(slide, spec.notes)
    return slide


def build_layout_b(prs: Presentation, spec: SlideSpec):
    """Build a section opener slide."""

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    if spec.mosaic_images:
        add_mosaic_panel(slide, spec.mosaic_images)
    else:
        add_image(
            slide,
            spec.image,
            0,
            0,
            int(SLIDE_WIDTH),
            int(SLIDE_HEIGHT),
            mode='cover',
        )
    add_rect(
        slide,
        0,
        int(SLIDE_HEIGHT) - int(SECTION_BAR_HEIGHT),
        int(SLIDE_WIDTH),
        int(SECTION_BAR_HEIGHT),
        BG,
        transparency=0.25,
    )
    add_textbox(
        slide,
        int(MARGIN),
        int(SLIDE_HEIGHT) - int(Cm(3.4)),
        Cm(2.2),
        Cm(1.6),
        spec.section_number,
        font_size=SECTION_NUMBER_SIZE,
        color=LIGHT_ACCENT,
        bold=True,
    )
    add_textbox(
        slide,
        int(MARGIN),
        int(SLIDE_HEIGHT) - int(Cm(2.2)),
        int(SLIDE_WIDTH) - int(Cm(4.0)),
        Cm(1.5),
        spec.title,
        font_size=Pt(30),
        color=BODY,
        bold=True,
    )
    return finalize_slide(slide, spec)


def build_layout_c(prs: Presentation, spec: SlideSpec):
    """Build a content slide with left image and right bullets."""

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, spec.title)
    image_rect = spec.image_rect or (
        int(MARGIN),
        int(TITLE_BAR_HEIGHT) + int(Cm(0.7)),
        int(Cm(16.6)),
        int(Cm(13.6)),
    )
    bullet_rect = spec.bullet_rect or (
        int(Cm(18.8)),
        int(TITLE_BAR_HEIGHT) + int(Cm(1.0)),
        int(Cm(13.6)),
        int(Cm(12.8)),
    )
    for image in spec.extra_images:
        add_image_placement(slide, image)
    image_x, image_y, image_w, image_h = image_rect
    if spec.image:
        add_image(
            slide,
            spec.image,
            image_x,
            image_y,
            image_w,
            image_h,
            mode=spec.body_image_mode,
            placeholder_text=spec.placeholder_text,
        )
    elif spec.placeholder_text:
        add_placeholder(
            slide,
            image_x,
            image_y,
            image_w,
            image_h,
            spec.placeholder_text,
        )
    add_bullets(
        slide,
        spec.bullets,
        *bullet_rect,
        font_size=spec.body_font_size or BODY_SIZE,
    )
    return finalize_slide(slide, spec)


def build_layout_d(prs: Presentation, spec: SlideSpec):
    """Build a full-bleed image slide with title and caption overlays."""

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_image(
        slide,
        spec.image,
        0,
        0,
        int(SLIDE_WIDTH),
        int(SLIDE_HEIGHT),
        mode='cover',
    )
    if spec.title:
        add_rect(
            slide,
            0,
            0,
            int(SLIDE_WIDTH),
            int(Cm(1.9)),
            BG,
            transparency=0.3,
        )
        add_textbox(
            slide,
            int(MARGIN),
            Cm(0.3),
            int(SLIDE_WIDTH) - int(Cm(4.2)),
            Cm(1.1),
            spec.title,
            font_size=Pt(28),
            color=BODY,
            bold=True,
        )
    if spec.image_caption:
        add_rect(
            slide,
            0,
            int(SLIDE_HEIGHT) - int(FULL_BLEED_BAR_HEIGHT),
            int(SLIDE_WIDTH),
            int(FULL_BLEED_BAR_HEIGHT),
            BG,
            transparency=0.25,
        )
        add_textbox(
            slide,
            int(MARGIN),
            int(SLIDE_HEIGHT) - int(Cm(1.45)),
            int(SLIDE_WIDTH) - int(Cm(4.0)),
            Cm(0.8),
            spec.image_caption,
            font_size=CAPTION_SIZE,
            color=SECONDARY,
            italic=True,
        )
    return finalize_slide(slide, spec)


def build_layout_e(prs: Presentation, spec: SlideSpec):
    """Build a bullets-left, diagram-right slide."""

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, spec.title)
    bullet_rect = spec.bullet_rect or (
        int(MARGIN),
        int(TITLE_BAR_HEIGHT) + int(Cm(1.0)),
        int(Cm(13.5)),
        int(Cm(12.6)),
    )
    image_rect = spec.image_rect or (
        int(Cm(15.0)),
        int(TITLE_BAR_HEIGHT) + int(Cm(0.7)),
        int(Cm(17.8)),
        int(Cm(13.6)),
    )
    add_bullets(
        slide,
        spec.bullets,
        *bullet_rect,
        font_size=spec.body_font_size or BODY_SIZE,
    )
    image_x, image_y, image_w, image_h = image_rect
    add_image(
        slide,
        spec.image,
        image_x,
        image_y,
        image_w,
        image_h,
        mode='contain',
    )
    return finalize_slide(slide, spec)


def build_layout_f(prs: Presentation, spec: SlideSpec):
    """Build a text-led slide with optional call-to-action line."""

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, spec.title)
    bullet_rect = spec.bullet_rect or (
        int(MARGIN),
        int(TITLE_BAR_HEIGHT) + int(Cm(1.2)),
        int(SLIDE_WIDTH) - int(Cm(2.0)),
        int(Cm(9.8)),
    )
    add_bullets(
        slide,
        spec.bullets,
        *bullet_rect,
        font_size=spec.body_font_size or BODY_SIZE,
    )
    if spec.cta:
        cta_box = slide.shapes.add_textbox(
            int(MARGIN),
            int(SLIDE_HEIGHT) - int(Cm(4.4)),
            int(SLIDE_WIDTH) - int(Cm(2.0)),
            int(Cm(2.2)),
        )
        frame = cta_box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        add_markup_paragraph(
            frame,
            spec.cta,
            color=AMBER,
            font_size=CTA_SIZE,
            align=PP_ALIGN.CENTER,
        )
    return finalize_slide(slide, spec)


SLIDES: tuple[SlideSpec, ...] = (
    SlideSpec(
        layout='A',
        title='Fox Hunting & ARDF — M&K 2026 Program',
        subtitle='Mike & Key Amateur Radio Club (K7LED)',
        presenter='Tom KE4HET  |  June 20, 2026',
        notes=(
            'Talking point: welcome the club, introduce yourself, and '
            'frame this as a 2026 program preview.',
        ),
    ),
    SlideSpec(
        layout='F',
        title="Today's Agenda",
        bullets=(
            'What is fox hunting / ARDF?',
            'Our equipment — foxes, antennas, and a budget radio tip',
            'Radio Camp build event — Thu Jun 26 at Fort Flagler',
            'ARDF demonstration at Field Day — Jun 28–29',
            '2026 hunt season — July · August · September',
            'How to get involved',
            'Q&A',
        ),
        notes=(
            'Talking point: read the roadmap so the audience knows where '
            'the program is headed.',
            'Read the agenda aloud. Flag: RT-910B rec is coming in §2, '
            'hunt dates are tentative, questions welcome throughout.',
        ),
    ),
    SlideSpec(
        layout='B',
        title='What Is Fox Hunting / ARDF?',
        section_number='1',
        image='fox-hunting-friends.jpg',
        notes=(
            'Talking point: open section one by showing fox hunting as '
            'hands-on, approachable, and family friendly.',
        ),
    ),
    SlideSpec(
        layout='C',
        title='The Concept',
        image='a-fox-foxhunting.jpg',
        bullets=(
            'A hidden transmitter, the **"fox"**, is placed in the '
            'field',
            'Participants find it using **Radio Direction Finding '
            '(RDF)**',
            'Used worldwide for sport, training, and fun',
            'Any license class can participate',
        ),
        notes=(
            'Talking point: define the fox, RDF, and the basic activity '
            'in one clear example.',
        ),
    ),
    SlideSpec(
        layout='C',
        title='Three Types of Fox Hunting',
        image='Car_with_Fox_Hunt_Antenna.jpeg',
        body_image_mode='contain',
        body_font_size=Pt(16),
        image_rect=(
            int(Cm(7.2)),
            int(Cm(8.2)),
            int(Cm(10.0)),
            int(Cm(8.2)),
        ),
        extra_images=(
            ImagePlacement(
                image='DXE-OAB-Fox-Chasing-Image_Kids-on-the-Hunt_Post-10-18.jpg',
                x=int(MARGIN),
                y=int(TITLE_BAR_HEIGHT) + int(Cm(0.7)),
                w=int(Cm(16.6)),
                h=int(Cm(13.6)),
                mode='contain',
                placeholder_text=(
                    'DXE-OAB-Fox-Chasing-Image_Kids-on-the-Hunt_Post-10-18.jpg'
                ),
            ),
        ),
        bullets=(
            '**On-foot (ARDF):** milliwatt fox · parks & trails · '
            'walk/run with directional antenna',
            '**Local mobile:** <5 W fox · neighborhood roads · take '
            'bearings from car, hunt a local area',
            '**Regional mobile:** 5+ W fox · wide-area adventure · '
            'multi-city drive, final on-foot leg',
            'Power level of the fox **defines the hunt type** — higher '
            '= bigger search area',
            'All three share the same core skill: read the signal, take '
            'a bearing, make a move',
        ),
        notes=(
            'Three types: on-foot (milliwatts), local mobile (<5 W), '
            'regional (5+ W)',
            'Power level of the fox defines the scale of the hunt',
            'All three share the same core skill — the gear and distance '
            'change, not the technique',
            'This program covers all three across the summer',
        ),
    ),
    SlideSpec(
        layout='E',
        title='How Direction Finding Works',
        image='patterns.jpg',
        bullets=(
            'A **Yagi antenna** has a sharp directional pattern',
            'Rotate until signal is **strongest** → that\'s the bearing',
            'Take multiple bearings → triangulate the fox',
            '**Step attenuator** reduces signal when you\'re close',
        ),
        notes=(
            'Talking point: explain the bearing-taking loop and why '
            'attenuation matters near the fox.',
        ),
    ),
    SlideSpec(
        layout='B',
        title='Our Equipment',
        section_number='2',
        image='Mike_and_Key_Fox_Beacons.jpg',
        notes=(
            'Talking point: transition from the concept to the specific '
            'club gear and entry path.',
        ),
    ),
    SlideSpec(
        layout='C',
        title='The Fox — Byonics MicroFox',
        image='Mike_and_Key_Fox_Beacons.jpg',
        body_image_mode='contain',
        body_font_size=Pt(16),
        bullets=(
            '**MF-15** and **MF-PC** — both programmed and tested',
            'Frequency: **147.42 MHz** simplex, callsign K7LED',
            'Cycle: 15 s on / 15 s off with CW ID',
            'Sub-1 W output — suited for **on-foot and local mobile** '
            'hunts',
            'Regional hunts need **5+ W** — 10 W Pi Zero beacon planned '
            '(issue #4)',
            'WWARA coordination in progress (Steve N9VW)',
            'Deploy **one at a time** (currently identical config)',
        ),
        notes=(
            'Talking point: describe the current beacon inventory and '
            'how each fox is configured for the program.',
            'Sub-1 W output means these work for on-foot and local '
            'mobile, not regional',
            'Regional range needs 5+ W — the Pi Zero beacon project '
            'addresses that',
        ),
    ),
    SlideSpec(
        layout='C',
        title="The Hunter's Antenna — Tape Measure Yagi",
        image='Tapemeasure Yagi.jpg',
        bullets=(
            '**3-element** tape-measure Yagi — 2 m band',
            'Lightweight, collapsible, inexpensive to build',
            'Connects to any HT with a BNC connector',
            "**We'll build these at Radio Camp** on Thu Jun 26",
        ),
        notes=(
            'Talking point: position the tape-measure Yagi as the '
            'practical, low-cost hunter antenna.',
        ),
    ),
    SlideSpec(
        layout='C',
        title='The Close-In Tool — Step Attenuator',
        image='rf_step_att_2.jpg',
        bullets=(
            'Inserted **between HT and antenna**',
            'Reduces signal strength in steps as you approach',
            'Prevents "saturated S-meter" confusion near the fox',
            '**Also built at Radio Camp**, Project 2',
        ),
        notes=(
            'Talking point: explain how attenuation keeps close-in '
            'hunting readable and less frustrating.',
        ),
    ),
    SlideSpec(
        layout='C',
        title='Entry-Level Radio — Radtel RT-910B',
        image='71PP0t4+U4L._AC_SL1500_.jpg',
        body_image_mode='contain',
        image_rect=(
            int(MARGIN),
            int(TITLE_BAR_HEIGHT) + int(Cm(0.7)),
            int(Cm(13.5)),
            int(Cm(13.6)),
        ),
        bullet_rect=(
            int(Cm(15.6)),
            int(TITLE_BAR_HEIGHT) + int(Cm(1.0)),
            int(Cm(16.7)),
            int(Cm(12.8)),
        ),
        bullets=(
            '~**$30** on Amazon — lowest barrier to entry',
            'Large **S-meter** display — great visual RDF aid',
            '**TX lockout** — safe for loaner / demo use',
            'BNC-to-SMA adapter required (~$5, 2-pack)',
            '⚠️ Conditional recommendation — not yet hands-on tested',
        ),
        notes=(
            'Talking point: share the lowest-cost HT path while keeping '
            'the recommendation clearly provisional.',
            'Emphasize this is conditional — we haven\'t put one in our '
            'hands yet.',
        ),
    ),
    SlideSpec(
        layout='B',
        title='Radio Camp Build Event',
        section_number='3',
        image='DIY 70cm handheld Yagi.png',
        notes=(
            'Talking point: open the build-event section with a strong '
            'maker image and clear date anchor.',
        ),
    ),
    SlideSpec(
        layout='C',
        title='Thu Jun 26 — Fort Flagler, Wagon Wheel',
        image='Tapemeasure_Yagi-2.jpg',
        body_image_mode='contain',
        bullets=(
            '**Radio Camp:** Mon Jun 22 – Mon Jun 29 at Fort Flagler '
            'State Park',
            'Build session **Thursday afternoon** (~5 hours, time TBD on '
            'site)',
            '**Project 1:** 3-element 2 m tape-measure Yagi *(required)*',
            '**Project 2:** Step attenuator *(if time allows)*',
            '**Project 3:** Dual-band 2m/70cm Yagi *(optional, advanced)*',
        ),
        notes=(
            'Talking point: outline the camp timing, venue, and three '
            'project tracks for the Thursday build.',
        ),
    ),
    SlideSpec(
        layout='F',
        title='Who, What, How',
        bullets=(
            'Licensed amateurs at camp — **1–10 participants**',
            'Hands-on from minute one — no whiteboard lectures',
            'Tom KE4HET brings materials for **at least 5 sets**',
            'Materials mostly on hand; everything expected by Thursday',
            '🙋 Who wants to build an antenna AND will be at Radio Camp on '
            'Thursday?',
        ),
        cta='🙋 Volunteers needed — build assistants & tool wranglers',
        notes=(
            'Talking point: make the build feel approachable and end '
            'with a specific volunteer ask.',
        ),
    ),
    SlideSpec(
        layout='B',
        title='ARDF Demonstration at Field Day',
        section_number='4',
        image='Ft_Flagler_Battery_Grattan.webp',
        notes=(
            'Talking point: shift from club-internal building to the '
            'public-facing Field Day demonstration.',
        ),
    ),
    SlideSpec(
        layout='E',
        title='Sat–Sun Jun 28–29 — Fort Flagler',
        image='Ft_Flagler_Google_Maps_Satelite.png',
        bullet_rect=(
            int(MARGIN),
            int(TITLE_BAR_HEIGHT) + int(Cm(1.0)),
            int(Cm(15.6)),
            int(Cm(12.6)),
        ),
        image_rect=(
            int(Cm(17.4)),
            int(TITLE_BAR_HEIGHT) + int(Cm(0.7)),
            int(Cm(15.6)),
            int(Cm(13.6)),
        ),
        bullets=(
            '**Visitor tent:** Battery Grattan (start point)',
            '**Fox:** hidden south side of Parade Grounds',
            'Fox location unstaffed — visible from tent',
            'Frequency: **147.42 MHz** (WWARA pending)',
        ),
        notes=(
            'Talking point: orient the audience to the demo geography and '
            'the simple start-to-finish hunt path.',
        ),
    ),
    SlideSpec(
        layout='C',
        title='How the Demo Works',
        image='Kids_with_tape-measure_yagi.webp',
        bullets=(
            'Option 1: Line-of-sight demonstration of directional antenna',
            'Option 2: Fox at the Wagon Wheel, guest borrows equipment and '
            'hunts their way to the Fox for a prize / certificate',
        ),
        notes=(
            'Talking point: walk through the visitor experience so club '
            'members can picture the staffing model.',
        ),
    ),
    SlideSpec(
        layout='F',
        title='We Need Your Help',
        bullets=(
            '🙋 **Demo hosts** — explain fox hunting to the public',
            '🦊 Hunter wranglers, return hunting gear to visitor tent',
            '📻 **Loaner gear** — HTs and Yagis for visitors',
        ),
        cta='Talk to Tom KE4HET; RSVP requested (not required)',
        notes=(
            'Talking point: make the staffing and loaner-equipment asks '
            'explicit for Field Day.',
        ),
    ),
    SlideSpec(
        layout='B',
        title='2026 Hunt Season',
        section_number='5',
        mosaic_images=(
            'Marymoor_Park_Google_Satilite.png',
            'Seattle_Neighborhood_Google_Satelite.png',
            'West_King_county_Google_Satelite.png',
        ),
        notes=(
            'Three-event progression spanning all three hunt types',
            'July = on-foot milliwatts, August = on-foot milliwatts, '
            'September = regional 5+ W',
            'All dates are tentative',
        ),
    ),
    SlideSpec(
        layout='C',
        title='July — Practice Hunt *(tentative)*',
        image='Flaming_Geyser_State_Park_Sign.webp',
        bullets=(
            '**Type: on-foot ARDF · milliwatt fox · parks & trails**',
            '**Sun Jul 19** — Flaming Geyser State Park, Auburn',
            'On-foot, single fox, noon – 2:30 PM (2.5 hrs)',
            '480+ acres, Green River valley, wooded trails',
            'Discover Pass required per vehicle; carpools encouraged',
            'Good intro event — multipath adds realistic challenge',
        ),
        notes=(
            'Talking point: present July as the entry-friendly practice '
            'hunt with realistic terrain and manageable logistics.',
        ),
    ),
    SlideSpec(
        layout='C',
        title='August — On-Foot Hunt *(tentative)*',
        image='Marymoor_Park_Google_Satilite.png',
        body_image_mode='contain',
        bullets=(
            '**Type: on-foot ARDF · milliwatt fox · parks & trails**',
            '**Sat Aug 22** — Marymoor Park, Redmond',
            'On-foot, single fox, noon – 2:30 PM',
            '640 acres — wooded Sammamish River corridor + open fields',
            'Light-rail accessible (2 Line)',
            'Note: wrap by 2:30 PM — Marymoor Live concert that evening',
        ),
        notes=(
            'Talking point: describe Marymoor as the main on-foot hunt, '
            'noting transit access and the strict afternoon wrap time.',
        ),
    ),
    SlideSpec(
        layout='C',
        title='September — Regional Mobile Hunt *(tentative)*',
        image='Car_with_Fox_Hunt_Antenna.jpeg',
        body_image_mode='contain',
        body_font_size=Pt(16),
        bullets=(
            '**Type: regional mobile · 5+ W fox · wide-area adventure**',
            '**Sun Sep 20** — Finish at Howarth Park, Everett',
            'Hunt from your **vehicle** — take bearings while driving',
            'Final leg: short **on-foot** walk on forested hillside',
            'Puget Sound views; social wrap-up at finish',
            '⚠️ Current MicroFox units sub-1 W — **too low for regional '
            'range**; 10 W Pi Zero beacon planned',
            '⚠️ Stay clear of BNSF railroad tracks at finish',
        ),
        notes=(
            'Talking point: position September as the mobile hunt, with '
            'a scenic finish and a clear safety reminder.',
        ),
    ),
    SlideSpec(
        layout='E',
        title="The Region We'll Hunt",
        image='Greater_Seattle_Google_satelite.png',
        bullets=(
            'Three events across the greater Puget Sound area',
            'July 19: **Auburn** (Flaming Geyser)',
            'August 22: **Redmond** (Marymoor)',
            'September 20: Regional hunt, details TBD',
        ),
        notes=(
            'Talking point: zoom out and show the three hunt areas as a '
            'regional program, not isolated events.',
        ),
    ),
    SlideSpec(
        layout='D',
        title='How to Get Involved',
        image='Group_of_fox_hunters.jpg',
        notes=(
            'Talking point: pivot from the event calendar to the specific '
            'ways members can jump in right away.',
        ),
    ),
    SlideSpec(
        layout='F',
        title='Join the Program',
        bullets=(
            '**Radio Camp (Jun 22–29):** attend and/or volunteer Thu '
            'build',
            '**Field Day demo (Jun 28–29):** host visitors or hide the '
            'fox',
            '**Hunts:** open participation; check in with net control at '
            'start or when you begin hunting',
            '**RSVP:** FoxHunting@mkarc.groups.io requested, not required',
            '🌐 https://mikeandkey.org/foxhunt.php',
            '**Questions:** Tom KE4HET',
        ),
        qr_image='qr_signup.png',
        qr_rect=(
            int(SLIDE_WIDTH) - int(MARGIN) - int(Cm(3.0)),
            int(SLIDE_HEIGHT) - int(MARGIN) - int(Cm(3.0)),
            int(Cm(3.0)),
            int(Cm(3.0)),
        ),
        show_logo=False,
        notes=(
            'Talking point: finish the call to action with the top entry '
            'points, contact path, and mailing list.',
        ),
    ),
    SlideSpec(
        layout='A',
        subtitle='Tom KE4HET  |  https://mikeandkey.org/foxhunt.php',
        qr_image='qr_signup.png',
        qr_rect=(
            (int(SLIDE_WIDTH) - int(Cm(4.0))) // 2,
            int(Cm(12.7)),
            int(Cm(4.0)),
            int(Cm(4.0)),
        ),
        notes=(
            'Talking point: invite questions and leave the audience with '
            'a clear follow-up contact path.',
        ),
    ),
)


def build_presentation() -> Presentation:
    """Create and populate the full presentation object."""

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    prs.core_properties.title = 'Fox Hunting & ARDF — M&K 2026 Program'
    prs.core_properties.author = 'GitHub Copilot'
    prs.core_properties.subject = 'Mike & Key ARC presentation'

    for index, spec in enumerate(SLIDES, start=1):
        if spec.layout == 'A':
            build_layout_a(prs, spec, closing=index == len(SLIDES))
        elif spec.layout == 'B':
            build_layout_b(prs, spec)
        elif spec.layout == 'C':
            build_layout_c(prs, spec)
        elif spec.layout == 'D':
            build_layout_d(prs, spec)
        elif spec.layout == 'E':
            build_layout_e(prs, spec)
        elif spec.layout == 'F':
            build_layout_f(prs, spec)
        else:
            raise ValueError(f'Unsupported layout: {spec.layout}')
    return prs


def main() -> None:
    """Build the presentation and write it to disk."""

    prs = build_presentation()
    prs.save(OUTPUT_PATH)
    verify = Presentation(str(OUTPUT_PATH))
    slide_count = len(verify.slides)
    if slide_count != len(SLIDES):
        raise RuntimeError(
            f'Expected {len(SLIDES)} slides, found {slide_count}.'
        )
    print(f'Created {OUTPUT_PATH} with {slide_count} slides.')


if __name__ == '__main__':
    main()
